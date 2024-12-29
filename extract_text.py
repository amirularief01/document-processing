from PyPDF2 import PdfReader
from pptx import Presentation
import tika
from tika import parser

def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_with_tika(file_path):
    """
    Extracts text from any document using Apache Tika.
    Supports PDFs, PPTX, DOCX, and more.
    """
    tika.initVM()
    parsed = parser.from_file(file_path)
    text = parsed.get("content", "").strip()  # Extract content and clean up
    return text